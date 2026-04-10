#!/usr/bin/env python3
"""
build_app_pptx.py

Reconstructs the 25 slides from workshop_app/src/App.tsx as an editable
.pptx file. White background, Harvard crimson accent, simple card layouts.
All text is real text (editable in PowerPoint/Keynote), not screenshots.

Run:
    /Users/cy/dev/hdsi/.venv/bin/python build_app_pptx.py

Output:
    /Users/cy/dev/hdsi/workshop_session4/slides_app.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Slide geometry (16:9) ────────────────────────────────────────
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)
BLANK = prs.slide_layouts[6]

# ── Colors ───────────────────────────────────────────────────────
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK    = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_BODY    = RGBColor(0x37, 0x41, 0x51)   # gray-700
TEXT_GRAY    = RGBColor(0x6B, 0x72, 0x80)   # gray-500
TEXT_LIGHT   = RGBColor(0x9C, 0xA3, 0xAF)   # gray-400
RED          = RGBColor(0x99, 0x1B, 0x1B)   # red-800 / Harvard crimson
RED_BRIGHT   = RGBColor(0xB9, 0x1C, 0x1C)   # red-700
RED_DARK     = RGBColor(0x7F, 0x1D, 0x1D)   # red-900
RED_TINT     = RGBColor(0xFE, 0xF2, 0xF2)   # red-50
EMERALD      = RGBColor(0x04, 0x78, 0x57)   # emerald-700
EMERALD_TINT = RGBColor(0xEC, 0xFD, 0xF5)   # emerald-50
ROSE         = RGBColor(0xBE, 0x12, 0x3C)   # rose-700
ROSE_TINT    = RGBColor(0xFF, 0xF1, 0xF2)   # rose-50
AMBER        = RGBColor(0xB4, 0x53, 0x09)   # amber-700
AMBER_TINT   = RGBColor(0xFF, 0xFB, 0xEB)   # amber-50
CYAN         = RGBColor(0x0E, 0x74, 0x90)   # cyan-700
CYAN_TINT    = RGBColor(0xEC, 0xFE, 0xFF)   # cyan-50
CARD_BG      = RGBColor(0xF9, 0xFA, 0xFB)   # gray-50
CARD_BORDER  = RGBColor(0xE5, 0xE7, 0xEB)   # gray-200
CODE_BG      = RGBColor(0xF3, 0xF4, 0xF6)   # gray-100

FONT_SANS = "Inter"
FONT_MONO = "Menlo"

# ── Helpers ──────────────────────────────────────────────────────
def new_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide

def add_text(slide, x, y, w, h, text, *, size=12, color=TEXT_DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_SANS):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box

def add_multiline(slide, x, y, w, h, lines, *, size=12, color=TEXT_BODY, bold=False, italic=False,
                  align=PP_ALIGN.LEFT, font=FONT_SANS, line_spacing=1.15):
    """lines: list of (text, dict) where dict can override size/color/bold/italic/font."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = opts.get("font", font)
        run.font.size = Pt(opts.get("size", size))
        run.font.color.rgb = opts.get("color", color)
        run.font.bold = opts.get("bold", bold)
        run.font.italic = opts.get("italic", italic)
    return box

def add_card(slide, x, y, w, h, *, fill=CARD_BG, border=CARD_BORDER, border_w=0.75, radius=0.08):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border
    rect.line.width = Pt(border_w)
    rect.adjustments[0] = radius
    rect.shadow.inherit = False
    return rect

def add_title(slide, text, *, size=32):
    add_text(slide, Inches(0.5), Inches(0.55), Inches(SLIDE_W_IN - 1), Inches(0.85),
             text, size=size, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

def add_subtitle(slide, text):
    add_text(slide, Inches(0.5), Inches(1.45), Inches(SLIDE_W_IN - 1), Inches(0.45),
             text, size=15, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
def slide_01_title():
    s = new_slide()
    add_text(s, Inches(0.5), Inches(2.7), Inches(SLIDE_W_IN - 1), Inches(0.9),
             "Coding and Doing Math Faster", size=44, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.55), Inches(SLIDE_W_IN - 1), Inches(0.9),
             "with LLMs and Agentic Workflows", size=44, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.7), Inches(SLIDE_W_IN - 1), Inches(0.5),
             "Generative AI for Scholarship — Session 4", size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "HARVARD DATA SCIENCE INITIATIVE", size=11, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Where We Are
# ════════════════════════════════════════════════════════════════
def slide_02_where_we_are():
    s = new_slide()
    add_title(s, "Where We Are")
    add_subtitle(s, "Four sessions. Today is the last one.")
    rows = [
        ("1", "The Basics", "Gemini, NotebookLM, Gems", "", False),
        ("2", "Hugging Face + MCP", "datasets, transformers, pipelines", "Audirac", False),
        ("3", "Claude Code CLI", "agentic coding from the terminal", "Stubbs", False),
        ("4", "From conversations to methodologies", "skills, agents, your methodology", "Today", True),
    ]
    y = Inches(2.2)
    for n, title, desc, instr, hl in rows:
        x = Inches(2.0); w = Inches(SLIDE_W_IN - 4.0); h = Inches(0.8)
        fill = RED_TINT if hl else CARD_BG
        border = RED_BRIGHT if hl else CARD_BORDER
        add_card(s, x, y, w, h, fill=fill, border=border, border_w=1.5 if hl else 0.75)
        # number badge
        badge_color = RED_BRIGHT if hl else TEXT_LIGHT
        add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(0.5), Inches(0.5),
                 n, size=22, color=badge_color, bold=True, align=PP_ALIGN.CENTER)
        # title + desc
        add_text(s, x + Inches(0.95), y + Inches(0.13), Inches(w.inches - 2.3), Inches(0.32),
                 title, size=14, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(0.95), y + Inches(0.45), Inches(w.inches - 2.3), Inches(0.3),
                 desc, size=11, color=TEXT_GRAY)
        # instructor
        if instr:
            add_text(s, x + Inches(w.inches - 1.4), y + Inches(0.27), Inches(1.2), Inches(0.3),
                     instr.upper(), size=10, color=TEXT_GRAY, bold=True, align=PP_ALIGN.RIGHT)
        y += Inches(0.95)
    add_text(s, Inches(0.5), Inches(6.95), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "Session 3 was a conversation. Today we go from conversations to methodologies.",
             size=11, color=TEXT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — AI Coding Landscape
# ════════════════════════════════════════════════════════════════
def slide_03_landscape():
    s = new_slide()
    add_title(s, "The AI Coding Landscape")
    add_subtitle(s, "Three layers of abstraction. Three different users.")
    cards = [
        ("Browser Tools", "Lovable, v0, Google AI Studio", "Everyone", ROSE),
        ("IDE Agents", "Cursor, VS Code + Copilot, Windsurf", "Professional Humans", CYAN),
        ("CLI Agents", "Claude Code, Codex CLI, Gemini CLI", "Humans + Agents", RED_BRIGHT),
    ]
    card_w = 3.6; gap = 0.4
    total = card_w * 3 + gap * 2
    start_x = (SLIDE_W_IN - total) / 2
    for i, (t, desc, user, accent) in enumerate(cards):
        x = Inches(start_x + i * (card_w + gap)); y = Inches(2.4); w = Inches(card_w); h = Inches(3.6)
        add_card(s, x, y, w, h)
        # accent bar at top
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        # title
        add_text(s, x + Inches(0.3), y + Inches(0.5), w - Inches(0.6), Inches(0.5),
                 t, size=20, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        # desc
        add_text(s, x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), Inches(0.8),
                 desc, size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
        # user pill at bottom
        pill_y = y + h - Inches(0.8)
        pill = add_card(s, x + Inches(0.6), pill_y, w - Inches(1.2), Inches(0.45),
                        fill=CARD_BG, border=CARD_BORDER, radius=0.4)
        add_text(s, x + Inches(0.6), pill_y + Inches(0.08), w - Inches(1.2), Inches(0.3),
                 user, size=10, color=TEXT_BODY, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Today's Roadmap
# ════════════════════════════════════════════════════════════════
def slide_04_roadmap():
    s = new_slide()
    add_title(s, "Today's Roadmap")
    items = [
        ("1", "First Principles", "What is an agent?"),
        ("2", "Skills: The Core", "What they are. Why they matter."),
        ("3", "Build Your Own", "Hands-on for your research"),
        ("4", "Synthesize", "Share, discuss, next steps"),
    ]
    card_w = 2.7; gap = 0.3
    total = card_w * 4 + gap * 3
    start_x = (SLIDE_W_IN - total) / 2
    for i, (n, title, desc) in enumerate(items):
        x = Inches(start_x + i * (card_w + gap)); y = Inches(2.8); w = Inches(card_w); h = Inches(2.3)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.4), y + Inches(0.5), w - Inches(0.8), Inches(0.7),
                 f"{n}. {title}", size=18, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(0.4), y + Inches(1.3), w - Inches(0.8), Inches(0.8),
                 desc, size=12, color=TEXT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Chatbot vs Agent
# ════════════════════════════════════════════════════════════════
def slide_05_chatbot_vs_agent():
    s = new_slide()
    add_title(s, "Chatbot vs. Agent")
    add_subtitle(s, "What makes an agent, really?")
    cols = [
        ("Chatbot", ["Reactive", "Single-turn", "No tools"], CARD_BG, CARD_BORDER, TEXT_GRAY),
        ("Agent", ["Proactive", "Multi-step", "Uses tools"], RED_TINT, RED_BRIGHT, TEXT_DARK),
    ]
    col_w = 4.2; gap = 0.6
    total = col_w * 2 + gap
    start_x = (SLIDE_W_IN - total) / 2
    for i, (t, items, fill, border, title_color) in enumerate(cols):
        x = Inches(start_x + i * (col_w + gap)); y = Inches(2.3); w = Inches(col_w); h = Inches(3.0)
        add_card(s, x, y, w, h, fill=fill, border=border, border_w=1.5 if i == 1 else 0.75)
        add_text(s, x + Inches(0.5), y + Inches(0.45), w - Inches(1), Inches(0.5),
                 t, size=24, color=title_color, bold=True)
        for j, it in enumerate(items):
            add_text(s, x + Inches(0.5), y + Inches(1.2 + j * 0.5), w - Inches(1), Inches(0.4),
                     "→  " + it, size=15, color=title_color)
    add_text(s, Inches(1), Inches(5.7), Inches(SLIDE_W_IN - 2), Inches(1.2),
             "An agent is a system that uses an LLM to pursue a goal by planning and "
             "taking actions with tools, iterating in a feedback loop until it's done.",
             size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — The Agent Loop (code block)
# ════════════════════════════════════════════════════════════════
def slide_06_agent_loop():
    s = new_slide()
    add_title(s, "The Agent Loop")
    add_subtitle(s, "It's just a while loop.")
    x, y, w, h = Inches(2.2), Inches(2.2), Inches(SLIDE_W_IN - 4.4), Inches(4.4)
    add_card(s, x, y, w, h, fill=CODE_BG, border=CARD_BORDER, radius=0.05)
    # filename header
    add_text(s, x + Inches(0.4), y + Inches(0.2), Inches(3), Inches(0.3),
             "agent_loop.py", size=10, color=TEXT_GRAY, font=FONT_MONO)
    # code lines
    code = [
        ("while not done:", {"color": RED_BRIGHT, "bold": True}),
        ("    # 1. Think: what should I do next?", {"color": TEXT_LIGHT, "italic": True}),
        ("    plan = llm.think(goal, memory, observations)", {"color": TEXT_DARK}),
        ("", {}),
        ("    # 2. Act: use a tool", {"color": TEXT_LIGHT, "italic": True}),
        ("    result = use_tool(plan.next_action())", {"color": TEXT_DARK}),
        ("", {}),
        ("    # 3. Observe: what happened?", {"color": TEXT_LIGHT, "italic": True}),
        ("    memory.add(result)", {"color": TEXT_DARK}),
        ("", {}),
        ("    # 4. Decide: am I done?", {"color": TEXT_LIGHT, "italic": True}),
        ("    done = plan.is_complete()", {"color": TEXT_DARK}),
    ]
    add_multiline(s, x + Inches(0.4), y + Inches(0.7), w - Inches(0.8), h - Inches(1),
                  code, size=14, font=FONT_MONO, color=TEXT_DARK, line_spacing=1.25)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Four Components
# ════════════════════════════════════════════════════════════════
def slide_07_four_components():
    s = new_slide()
    add_title(s, "The Four Components of an Agent")
    add_subtitle(s, "Each one swappable. Each one essential.")
    items = [
        ("LLM", "the brain", "Reasoning, planning, deciding what to do next.", RED_BRIGHT),
        ("Tools", "the hands", "Read files, write code, run scripts, call APIs.", CYAN),
        ("Memory", "the notebook", "Conversation history, project context, learned facts.", EMERALD),
        ("Planning", "the strategy", "Breaking goals into steps, sequencing actions.", AMBER),
    ]
    card_w = 5.8; card_h = 2.0; gap = 0.3
    start_x = (SLIDE_W_IN - card_w * 2 - gap) / 2
    for i, (label, role, desc, accent) in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(start_x + col * (card_w + gap))
        y = Inches(2.4 + row * (card_h + gap))
        w = Inches(card_w); h = Inches(card_h)
        add_card(s, x, y, w, h)
        # ChevronRight ">" marker
        add_text(s, x + Inches(0.4), y + Inches(0.4), Inches(0.4), Inches(0.6),
                 "›", size=32, color=accent, bold=True)
        # label
        add_text(s, x + Inches(0.95), y + Inches(0.35), Inches(card_w - 1.2), Inches(0.5),
                 label, size=22, color=TEXT_DARK, bold=True)
        # role (italic subtitle)
        add_text(s, x + Inches(0.95), y + Inches(0.85), Inches(card_w - 1.2), Inches(0.4),
                 role, size=14, color=TEXT_GRAY, italic=True)
        # desc
        add_text(s, x + Inches(0.4), y + Inches(1.35), Inches(card_w - 0.8), Inches(0.6),
                 desc, size=12, color=TEXT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Workflows vs Agents
# ════════════════════════════════════════════════════════════════
def slide_08_workflows_vs_agents():
    s = new_slide()
    add_title(s, "Workflows vs. Agents")
    add_subtitle(s, "Most research tasks need something in between.")
    cols = [
        ("Workflow", ["Predefined code paths", "Deterministic", "You write the control flow", "Like a data pipeline", "Reliable but rigid"]),
        ("Agent", ["LLM decides what's next", "Dynamic, adaptive", "Agent chooses the path", "Like a research assistant", "Flexible but undirected"]),
    ]
    col_w = 5.5; gap = 0.5
    total = col_w * 2 + gap
    start_x = (SLIDE_W_IN - total) / 2
    for i, (t, items) in enumerate(cols):
        x = Inches(start_x + i * (col_w + gap)); y = Inches(2.3); w = Inches(col_w); h = Inches(3.4)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.5), y + Inches(0.4), w - Inches(1), Inches(0.5),
                 t, size=22, color=TEXT_DARK, bold=True)
        for j, it in enumerate(items):
            add_text(s, x + Inches(0.5), y + Inches(1.1 + j * 0.42), w - Inches(1), Inches(0.4),
                     "→  " + it, size=12, color=TEXT_BODY)
    # Bottom highlight box
    bx = Inches((SLIDE_W_IN - 9) / 2); by = Inches(6.0); bw = Inches(9); bh = Inches(0.95)
    add_card(s, bx, by, bw, bh, fill=RED_TINT, border=RED_BRIGHT, border_w=1.0)
    add_text(s, bx, by + Inches(0.15), bw, Inches(0.35),
             "Skills are the something in between.", size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, bx, by + Inches(0.5), bw, Inches(0.4),
             "Pure workflows can't adapt. Pure agents won't impose discipline. Skills fix the structure and let the agent fill in the content.",
             size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Claude Code IS an Agent
# ════════════════════════════════════════════════════════════════
def slide_09_claude_code_is():
    s = new_slide()
    add_title(s, "Claude Code IS an Agent")
    add_subtitle(s, "The four components, instantiated.")
    rows = [
        ("LLM", "Claude Sonnet, Opus"),
        ("TOOLS", "Read · Write · Edit · Bash · Glob · Grep · WebFetch · Agent"),
        ("MEMORY", "Conversation context + CLAUDE.md + auto-memory"),
        ("PLANNING", "Plan mode (Shift+Tab), task tracking"),
    ]
    y = Inches(2.2)
    for label, value in rows:
        x = Inches(1.5); w = Inches(SLIDE_W_IN - 3); h = Inches(0.7)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.3), y + Inches(0.18), Inches(2), Inches(0.4),
                 label, size=18, color=RED, bold=True)
        add_text(s, x + Inches(2.5), y + Inches(0.2), w - Inches(2.8), Inches(0.4),
                 value, size=12, color=TEXT_DARK, font=FONT_MONO)
        y += Inches(0.85)
    # Subagents callout
    cy = Inches(5.85); cx = Inches(1.5); cw = Inches(SLIDE_W_IN - 3); ch = Inches(1.1)
    add_card(s, cx, cy, cw, ch, fill=CYAN_TINT, border=CYAN, border_w=1.0)
    add_text(s, cx + Inches(0.3), cy + Inches(0.15), Inches(3), Inches(0.4),
             "Subagents", size=14, color=CYAN, bold=True)
    add_text(s, cx + Inches(0.3), cy + Inches(0.55), cw - Inches(0.6), Inches(0.45),
             "Isolated sub-agent loops with their own context. Enable delegation, parallelism, sandboxing.",
             size=11, color=TEXT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — The Problem Agents Create
# ════════════════════════════════════════════════════════════════
def slide_10_problem():
    s = new_slide()
    add_title(s, "The Problem Agents Create")
    add_subtitle(s, "They can do anything, but they don't know YOUR methodology.")
    items = [
        ("Bayesian model", "Skips prior predictive checks"),
        ("Causal inference", "Jumps straight to estimation without a DAG"),
        ("Data analysis", "Picks arbitrary methods, forgets diagnostics"),
        ("Literature review", "Misses systematic checks, fabricates citations"),
    ]
    y = Inches(2.2)
    for ask, fail in items:
        x = Inches(2); w = Inches(SLIDE_W_IN - 4); h = Inches(0.7)
        add_card(s, x, y, w, h, fill=ROSE_TINT, border=ROSE, border_w=0.75)
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(4), Inches(0.4),
                 "Ask it to do " + ask, size=14, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(4.5), y + Inches(0.2), w - Inches(4.8), Inches(0.4),
                 "→  " + fail, size=12, color=ROSE)
        y += Inches(0.85)
    # Closing message
    cx = Inches(1.5); cy = Inches(5.85); cw = Inches(SLIDE_W_IN - 3); ch = Inches(1.2)
    add_card(s, cx, cy, cw, ch)
    add_text(s, cx, cy + Inches(0.25), cw, Inches(0.5),
             "Agents are like brilliant but untrained research assistants.",
             size=16, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx, cy + Inches(0.75), cw, Inches(0.4),
             "This is the problem skills solve.",
             size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — What Is a Skill (5-step chain)
# ════════════════════════════════════════════════════════════════
def slide_11_what_is_skill():
    s = new_slide()
    add_title(s, "What Is a Skill?")
    add_subtitle(s, "A structured instruction manual that teaches an agent HOW to do a specific task the right way.")
    steps = [
        ("One-off prompt", "one-shot", False),
        ("Custom command", "reusable text", False),
        ("CLAUDE.md", "always active", False),
        ("SKILL", "on-demand workflow", True),
        ("Subagent", "delegated isolated", False),
    ]
    cw = 2.3; gap = 0.15
    total = cw * 5 + gap * 4
    start_x = (SLIDE_W_IN - total) / 2
    for i, (name, desc, hl) in enumerate(steps):
        x = Inches(start_x + i * (cw + gap)); y = Inches(3.0); w = Inches(cw); h = Inches(1.8)
        fill = RED_TINT if hl else CARD_BG
        border = RED_BRIGHT if hl else CARD_BORDER
        bw = 2.0 if hl else 0.75
        add_card(s, x, y, w, h, fill=fill, border=border, border_w=bw)
        add_text(s, x, y + Inches(0.45), w, Inches(0.5),
                 name, size=14, color=RED if hl else TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.95), w, Inches(0.4),
                 desc, size=11, color=RED_BRIGHT if hl else TEXT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Why Skills Matter
# ════════════════════════════════════════════════════════════════
def slide_12_why_skills_matter():
    s = new_slide()
    add_title(s, "Why Skills Matter for Research")
    add_subtitle(s, "The before-and-after, item by item.")
    cols = [
        ("Without skills", [
            "Agent guesses your methodology",
            "No validation gates",
            "Different results each time",
            "Generic, surface-level output",
            "Skips what it doesn't know",
        ], ROSE_TINT, ROSE),
        ("With skills", [
            "Agent follows YOUR methodology",
            "Mandatory checkpoints",
            "Consistent, reproducible workflow",
            "Domain-appropriate output",
            "Encodes community knowledge",
        ], EMERALD_TINT, EMERALD),
    ]
    col_w = 5.5; gap = 0.4
    total = col_w * 2 + gap
    start_x = (SLIDE_W_IN - total) / 2
    for i, (t, items, fill, color) in enumerate(cols):
        x = Inches(start_x + i * (col_w + gap)); y = Inches(2.3); w = Inches(col_w); h = Inches(3.6)
        add_card(s, x, y, w, h, fill=fill, border=color, border_w=1.0)
        add_text(s, x + Inches(0.5), y + Inches(0.4), w - Inches(1), Inches(0.5),
                 t, size=18, color=color, bold=True)
        for j, it in enumerate(items):
            add_text(s, x + Inches(0.5), y + Inches(1.1 + j * 0.45), w - Inches(1), Inches(0.4),
                     "•  " + it, size=12, color=TEXT_BODY)
    add_text(s, Inches(0.5), Inches(6.3), Inches(SLIDE_W_IN - 1), Inches(0.5),
             '"Contains critical guardrails that agents won\'t apply unprompted."',
             size=11, color=TEXT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Skills vs Everything Else
# ════════════════════════════════════════════════════════════════
def slide_13_skills_vs_else():
    s = new_slide()
    add_title(s, "Skills vs. Everything Else")
    add_subtitle(s, "Adjacent concepts. Distinct purposes.")
    rows = [
        ("Prompts", "Reusable file on disk vs. one-off text you type", False),
        ("CLAUDE.md", "Invoked on demand vs. always active (lab protocol vs. dress code)", False),
        ("MCP", "Workflow instructions vs. tool connectivity (assembly manual vs. screwdriver)", True),
        ("Subagents", "Methodology guidance vs. task delegation", False),
    ]
    y = Inches(2.2)
    for right, contrast, bridge in rows:
        x = Inches(1.5); w = Inches(SLIDE_W_IN - 3); h = Inches(1.0)
        fill = CYAN_TINT if bridge else CARD_BG
        border = CYAN if bridge else CARD_BORDER
        add_card(s, x, y, w, h, fill=fill, border=border, border_w=1.0 if bridge else 0.75)
        add_text(s, x + Inches(0.4), y + Inches(0.15), Inches(4), Inches(0.4),
                 "Skills  ⚖  " + right, size=14, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(0.4), y + Inches(0.55), w - Inches(0.8), Inches(0.4),
                 contrast, size=11, color=TEXT_BODY)
        if bridge:
            add_text(s, x + Inches(0.4), y + Inches(0.78), w - Inches(0.8), Inches(0.25),
                     "↑ From Audirac's Session 2 you got the screwdrivers. Today you get the assembly manual.",
                     size=9, color=CYAN, italic=True)
        y += Inches(1.1)


# ════════════════════════════════════════════════════════════════
# SLIDE 14 — SKILL.md Anatomy
# ════════════════════════════════════════════════════════════════
def slide_14_skill_md_anatomy():
    s = new_slide()
    add_title(s, "Anatomy of a SKILL.md File")
    add_subtitle(s, "A methodology encoded as markdown.")
    # Code block
    cx, cy, cw, ch = Inches(1.0), Inches(2.2), Inches(8.5), Inches(4.6)
    add_card(s, cx, cy, cw, ch, fill=CODE_BG, border=CARD_BORDER)
    # File label
    add_text(s, cx + Inches(0.4), cy + Inches(0.2), Inches(3), Inches(0.3),
             "SKILL.md", size=10, color=TEXT_GRAY, font=FONT_MONO)
    code = [
        ("---", {"color": TEXT_LIGHT}),
        ("name: my-research-skill", {"color": TEXT_DARK}),
        ("description: >", {"color": TEXT_DARK}),
        ("  Opinionated workflow for X.", {"color": EMERALD}),
        ("  Use when asked to do Y or Z.", {"color": EMERALD}),
        ("---", {"color": TEXT_LIGHT}),
        ("", {}),
        ("# Skill Name", {"color": RED, "bold": True}),
        ("", {}),
        ("## Workflow overview", {"color": RED, "bold": True}),
        ("1. **Think** -- Define the question", {"color": TEXT_DARK}),
        ("2. **Validate** -- Check assumptions   ⚠️ ASK USER", {"color": AMBER}),
        ("3. **Implement** -- Write the code", {"color": TEXT_DARK}),
        ("", {}),
        ("## Critical rules", {"color": RED, "bold": True}),
        ("- Always do X before Y", {"color": TEXT_DARK}),
        ("- Never skip step 2", {"color": TEXT_DARK}),
    ]
    add_multiline(s, cx + Inches(0.4), cy + Inches(0.7), cw - Inches(0.8), ch - Inches(1),
                  code, size=11, font=FONT_MONO, line_spacing=1.2)
    # Side notes
    nx = Inches(SLIDE_W_IN - 2.7)
    notes = [
        ("Frontmatter", "Identifier & trigger phrases"),
        ("Workflow", "Sequential steps & checkpoints"),
        ("Critical Rules", "Iron laws & guardrails"),
    ]
    # Note: side notes overlap the code box on the right; instead place them below as a footer
    fy = Inches(6.95)
    for i, (lab, desc) in enumerate(notes):
        x = Inches(0.7 + i * 4.3)
        add_text(s, x, fy, Inches(2), Inches(0.3),
                 lab, size=10, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(1.4), fy, Inches(2.8), Inches(0.3),
                 desc, size=9, color=TEXT_GRAY)


# ════════════════════════════════════════════════════════════════
# SLIDE 15 — Bayesian Workflow Case Study (3 phases)
# ════════════════════════════════════════════════════════════════
def slide_15_bayesian_case():
    s = new_slide()
    add_title(s, "Case Study: Bayesian Workflow")
    add_subtitle(s, "Ten steps in three phases. Build → Run → Criticize & Report.")
    phases = [
        ("BUILD", CYAN, [
            ("1", "Generative story", "How were the data generated, in prose, before code"),
            ("2", "Specify priors", "With documented rationale, not 'Normal(0,1) and vibes'"),
            ("3", "Implement in PyMC", "coords, dims, nutpie sampler"),
        ]),
        ("RUN", AMBER, [
            ("4", "Prior predictive checks ⚠️", "Mandatory before sampling — the #1 violation"),
            ("5", "Inference", "nutpie, reproducible seeds from analysis name"),
            ("6", "Convergence diagnostics", "R̂ < 1.01, ESS, divergences, E-BFMI"),
        ]),
        ("CRITICIZE & REPORT", EMERALD, [
            ("7", "Posterior predictive", "Model criticism + LOO-PIT calibration"),
            ("8", "Prior sensitivity", "psense_summary on the InferenceData"),
            ("9", "Model comparison", "LOO-CV, ELPD with SE, prefer simpler when diff ≈ SE"),
            ("10", "Report", "94% HDI, no frequentist language, companion markdown"),
        ]),
    ]
    col_w = 4.1; gap = 0.25
    total = col_w * 3 + gap * 2
    start_x = (SLIDE_W_IN - total) / 2
    for i, (label, color, steps) in enumerate(phases):
        x = Inches(start_x + i * (col_w + gap)); y = Inches(2.2); w = Inches(col_w); h = Inches(4.7)
        add_card(s, x, y, w, h)
        # accent bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.07))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # phase label
        add_text(s, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.4),
                 label, size=11, color=color, bold=True)
        # steps
        sy = y + Inches(0.75)
        for n, name, desc in steps:
            add_text(s, x + Inches(0.3), sy, Inches(0.4), Inches(0.3),
                     n + ".", size=10, color=color, bold=True, font=FONT_MONO)
            add_text(s, x + Inches(0.65), sy, w - Inches(0.95), Inches(0.3),
                     name, size=12, color=TEXT_DARK, bold=True)
            add_text(s, x + Inches(0.65), sy + Inches(0.3), w - Inches(0.95), Inches(0.4),
                     desc, size=9, color=TEXT_GRAY)
            sy += Inches(0.78)
    add_text(s, Inches(0.5), Inches(7.05), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "Step 4 (prior predictive checks) is the step unscaffolded agents skip most often. The skill makes it mandatory.",
             size=10, color=TEXT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 16 — Bayesian Workflow Results
# ════════════════════════════════════════════════════════════════
def slide_16_bayesian_results():
    s = new_slide()
    add_title(s, "Bayesian Workflow — Results")
    add_subtitle(s, "Reproducible empirical claims, verified against benchmark.json")
    stats = [
        ("100%", "with skill", "53/53 individual checks", EMERALD, EMERALD_TINT),
        ("90.5%", "without skill", "48/53 individual checks", ROSE, ROSE_TINT),
        ("+29%", "wall-clock time", "the cost of rigor", AMBER, AMBER_TINT),
        ("+87%", "token count", "paid once per analysis", AMBER, AMBER_TINT),
    ]
    cw = 2.7; gap = 0.3
    total = cw * 4 + gap * 3
    start_x = (SLIDE_W_IN - total) / 2
    for i, (n, l, sub, color, tint) in enumerate(stats):
        x = Inches(start_x + i * (cw + gap)); y = Inches(2.5); w = Inches(cw); h = Inches(2.6)
        add_card(s, x, y, w, h, fill=tint, border=color, border_w=1.0)
        add_text(s, x, y + Inches(0.55), w, Inches(1.0),
                 n, size=46, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.65), w, Inches(0.4),
                 l, size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(2.05), w, Inches(0.4),
                 sub, size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    # Bottom prose
    add_text(s, Inches(1), Inches(5.7), Inches(SLIDE_W_IN - 2), Inches(0.5),
             "Six scenarios. 53 of 53 checks pass with the skill; 48 of 53 without. "
             "Without-skill range across scenarios: 88% – 100%.",
             size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(6.6), Inches(SLIDE_W_IN - 2), Inches(0.4),
             "baygent-skills/evals/bayesian-workflow/iteration-1/benchmark.json",
             size=9, color=TEXT_GRAY, font=FONT_MONO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 17 — Stress Test (with quote as title)
# ════════════════════════════════════════════════════════════════
def slide_17_stress_test():
    s = new_slide()
    add_text(s, Inches(0.5), Inches(0.4), Inches(SLIDE_W_IN - 1), Inches(0.85),
             '"Evals are necessary. They\'re not sufficient."',
             size=30, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(1.3), Inches(SLIDE_W_IN - 1), Inches(0.5),
             "Synthetic benchmarks catch what's predictable. Real research catches the rest.",
             size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    # Framing paragraph
    add_text(s, Inches(1.5), Inches(2.0), Inches(SLIDE_W_IN - 3), Inches(0.6),
             "The synthetic benchmark passed at v1.0. Then Andorra applied the skill to his own research — "
             "the Soccer Factor Model, an ordered logistic hierarchical model across 5 European leagues, "
             "~10,000 matches, 30+ features. Three bugs surfaced that the benchmark had not caught.",
             size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER, italic=True)
    # Three bug cards
    bugs = [
        ("1", "Median ≠ Mean",
         "Skill computed predictions with the posterior median.",
         "Medians don't commute with nonlinear transforms — the median of probabilities is not a probability.",
         "v1.1 critical rule: always use the posterior mean for predictive probabilities."),
        ("2", "No Sparsity Priors",
         "Skill had no guidance for 30+ feature regression.",
         "Default priors don't shrink enough; coefficients overfit and become uninterpretable.",
         "v1.1 added regularized (Finnish) horseshoe priors. Cites Piironen & Vehtari (2017)."),
        ("3", "Aggregate ≠ Per-Class",
         "Diagnostics passed at the aggregate level but missed per-class miscalibration.",
         "Ordered outcomes can be jointly calibrated yet wrong on individual classes.",
         "v1.1 added per-class calibration plots and Expected Calibration Error (ECE)."),
    ]
    cw = 4.0; gap = 0.2
    total = cw * 3 + gap * 2
    start_x = (SLIDE_W_IN - total) / 2
    for i, (n, t, broke, why, fix) in enumerate(bugs):
        x = Inches(start_x + i * (cw + gap)); y = Inches(3.4); w = Inches(cw); h = Inches(3.7)
        add_card(s, x, y, w, h)
        # number badge
        add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5),
                 n, size=18, color=ROSE, bold=True, align=PP_ALIGN.CENTER)
        # title
        add_text(s, x + Inches(0.85), y + Inches(0.27), w - Inches(1), Inches(0.4),
                 t, size=14, color=TEXT_DARK, bold=True)
        # WHAT BROKE
        add_text(s, x + Inches(0.3), y + Inches(0.85), w - Inches(0.6), Inches(0.25),
                 "WHAT BROKE", size=8, color=ROSE, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(1.05), w - Inches(0.6), Inches(0.6),
                 broke, size=10, color=TEXT_DARK)
        # WHY IT MATTERS
        add_text(s, x + Inches(0.3), y + Inches(1.7), w - Inches(0.6), Inches(0.25),
                 "WHY IT MATTERS", size=8, color=ROSE, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(1.9), w - Inches(0.6), Inches(0.7),
                 why, size=10, color=TEXT_GRAY)
        # FIX IN V1.1
        add_text(s, x + Inches(0.3), y + Inches(2.7), w - Inches(0.6), Inches(0.25),
                 "FIX IN V1.1", size=8, color=EMERALD, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(2.9), w - Inches(0.6), Inches(0.7),
                 fix, size=10, color=EMERALD)


# ════════════════════════════════════════════════════════════════
# SLIDE 18 — Causal Inference
# ════════════════════════════════════════════════════════════════
def slide_18_causal_inference():
    s = new_slide()
    add_title(s, "Case Study: Causal Inference")
    add_subtitle(s, "The thinking-first design.")
    cols = [
        ("Thinking Phase  (no code)", CYAN, [
            ("1", "Formulate causal question", False),
            ("2", "Draw the DAG", True),
            ("3", "Identify strategy", True),
            ("4", "Choose design", True),
        ]),
        ("Doing Phase", EMERALD, [
            ("5", "Estimate model", False),
            ("6", "Mandatory refutation tests", False),
            ("7", "Interpret with intervals", False),
            ("8", "Report assumptions first", False),
        ]),
    ]
    cw = 5.5; gap = 0.4
    total = cw * 2 + gap
    start_x = (SLIDE_W_IN - total) / 2
    for i, (label, color, steps) in enumerate(cols):
        x = Inches(start_x + i * (cw + gap)); y = Inches(2.3); w = Inches(cw); h = Inches(3.6)
        add_card(s, x, y, w, h)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.07))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, x + Inches(0.4), y + Inches(0.3), w - Inches(0.8), Inches(0.5),
                 label, size=18, color=TEXT_DARK, bold=True)
        sy = y + Inches(0.95)
        for n, name, confirm in steps:
            add_text(s, x + Inches(0.4), sy, Inches(0.4), Inches(0.4),
                     n + ".", size=12, color=color, bold=True, font=FONT_MONO)
            add_text(s, x + Inches(0.8), sy, w - Inches(2.5), Inches(0.4),
                     name, size=13, color=TEXT_DARK)
            if confirm:
                add_text(s, x + w - Inches(1.6), sy, Inches(1.4), Inches(0.4),
                         "⚠️ CONFIRM", size=9, color=AMBER, bold=True, align=PP_ALIGN.RIGHT)
            sy += Inches(0.55)
    add_text(s, Inches(0.5), Inches(6.3), Inches(SLIDE_W_IN - 1), Inches(0.4),
             'Key rule: "No estimation without a confirmed DAG."',
             size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.75), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "The 32-point benchmark gap is entirely reasoning discipline.",
             size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 19 — More Skills in the Wild
# ════════════════════════════════════════════════════════════════
def slide_19_more_skills():
    s = new_slide()
    add_title(s, "More Skills in the Wild")
    add_subtitle(s, "Three points on a spectrum. None of them about statistics.")
    cards = [
        ("Test-Driven Development", "github.com/obra/superpowers · Jesse Vincent",
         "RED → GREEN → REFACTOR. If the agent writes code before a failing test, the skill instructs it to delete the code and start over.",
         "Encodes a 20-year-old discipline that expert programmers know but routinely skip under deadline pressure. The skill removes the option to skip it.",
         EMERALD),
        ("Math Olympiad", "Anthropic official Claude Code plugin",
         "400+ lines orchestrating 8–12 parallel solver subagents · fresh verifiers blind to solver reasoning · asymmetric voting (4 to confirm, 2 to refute) · revision cycles",
         "Shows skills scale all the way to multi-agent orchestration. Calibrated abstention lets the agent refuse instead of confabulating — the same discipline as 'no estimation without a DAG.'",
         AMBER),
        ("Cross-Platform Spec", "agentskills.io · open standard · MIT",
         "The same file you write today runs unmodified in Claude Code, Cursor, Kimi Code, and Gemini CLI. baygent-skills, superpowers, and Anthropic's official plugins all conform.",
         "Skills are becoming infrastructure — like Python packages, but for methodology. Forkable, versionable, composable, openly licensed.",
         CYAN),
    ]
    cw = 4.0; gap = 0.2
    total = cw * 3 + gap * 2
    start_x = (SLIDE_W_IN - total) / 2
    for i, (t, attr, desc, why, color) in enumerate(cards):
        x = Inches(start_x + i * (cw + gap)); y = Inches(2.3); w = Inches(cw); h = Inches(4.8)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.4), y + Inches(0.4), w - Inches(0.8), Inches(0.45),
                 t, size=15, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(0.4), y + Inches(0.85), w - Inches(0.8), Inches(0.3),
                 attr, size=8, color=TEXT_GRAY, font=FONT_MONO)
        add_text(s, x + Inches(0.4), y + Inches(1.4), w - Inches(0.8), Inches(1.6),
                 desc, size=10, color=TEXT_BODY)
        # divider
        ly = y + Inches(3.3)
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.4), ly, w - Inches(0.8), Inches(0.01))
        line.fill.solid(); line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()
        add_text(s, x + Inches(0.4), ly + Inches(0.1), w - Inches(0.8), Inches(0.25),
                 "WHY IT MATTERS", size=8, color=color, bold=True)
        add_text(s, x + Inches(0.4), ly + Inches(0.35), w - Inches(0.8), Inches(1.0),
                 why, size=10, color=TEXT_BODY)


# ════════════════════════════════════════════════════════════════
# SLIDE 20 — Practical Example baygent-skills
# ════════════════════════════════════════════════════════════════
def slide_20_practical_example():
    s = new_slide()
    add_title(s, "The Practical Example")
    add_subtitle(s, "baygent-skills · MIT licensed · open it, read it, fork it.")
    # Code-block-style directory tree
    cx, cy, cw, ch = Inches(2.5), Inches(2.2), Inches(8.3), Inches(4.2)
    add_card(s, cx, cy, cw, ch, fill=CODE_BG, border=CARD_BORDER)
    add_text(s, cx + Inches(0.4), cy + Inches(0.2), Inches(3), Inches(0.3),
             "bayesian-workflow/", size=10, color=RED, font=FONT_MONO, bold=True)
    tree = [
        ("├──", "SKILL.md", "# 14KB · v1.2 · Alex Andorra"),
        ("├──", "README.md", ""),
        ("├──", "main.py", ""),
        ("├──", "pyproject.toml", ""),
        ("├──", "references/", "# progressive disclosure"),
        ("│   ├──", "priors.md", ""),
        ("│   ├──", "diagnostics.md", ""),
        ("│   ├──", "hierarchical.md", ""),
        ("│   ├──", "model-criticism.md", ""),
        ("│   ├──", "model-comparison.md", ""),
        ("│   ├──", "sensitivity.md", ""),
        ("│   └──", "reporting.md", ""),
        ("└──", "scripts/", "# executable diagnostics"),
        ("    ├──", "diagnose_model.py", ""),
        ("    └──", "calibration_check.py", ""),
    ]
    ty = cy + Inches(0.6)
    for prefix, name, comment in tree:
        line_text = f"{prefix} {name}"
        if comment:
            line_text += f"   {comment}"
        add_text(s, cx + Inches(0.4), ty, cw - Inches(0.8), Inches(0.22),
                 line_text, size=10, color=TEXT_DARK, font=FONT_MONO)
        ty += Inches(0.22)
    # clone command
    add_text(s, Inches(0.5), Inches(6.5), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "$ git clone https://github.com/Learning-Bayesian-Statistics/baygent-skills",
             size=11, color=EMERALD, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.95), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "Browsing this repo > any template.",
             size=11, color=TEXT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 21 — Thinking About YOUR Research
# ════════════════════════════════════════════════════════════════
def slide_21_your_research():
    s = new_slide()
    add_title(s, "Thinking About YOUR Research")
    add_subtitle(s, "Four questions to elicit your skill.")
    items = [
        ("1", "What task do you do repeatedly?", "data cleaning, model fitting, lit review, reviewing a paper…", False),
        ("2", "What does your advisor always remind you about?", "→ those are your critical rules", False),
        ("3", "What mistakes have you wished you'd been warned about?", "→ those are your common gotchas", False),
        ("4", "What would a new lab member need to know?", "→ that IS your skill", True),
    ]
    y = Inches(2.3)
    for n, q, a, hl in items:
        x = Inches(1.5); w = Inches(SLIDE_W_IN - 3); h = Inches(0.95)
        fill = RED_TINT if hl else CARD_BG
        border = RED_BRIGHT if hl else CARD_BORDER
        bw = 1.5 if hl else 0.75
        add_card(s, x, y, w, h, fill=fill, border=border, border_w=bw)
        add_text(s, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5),
                 n, size=18, color=RED if hl else TEXT_GRAY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.95), y + Inches(0.18), w - Inches(1.3), Inches(0.4),
                 q, size=14, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(0.95), y + Inches(0.55), w - Inches(1.3), Inches(0.4),
                 a, size=11, color=RED if hl else TEXT_GRAY, italic=True)
        y += Inches(1.05)


# ════════════════════════════════════════════════════════════════
# SLIDE 22 — Hands-On Exercise
# ════════════════════════════════════════════════════════════════
def slide_22_hands_on():
    s = new_slide()
    add_title(s, "Hands-On Exercise")
    add_subtitle(s, "Two tracks. Pick one.")
    cols = [
        ("TRACK A", "Write your own SKILL.md", [
            "Choose a workflow from your research",
            "Write SKILL.md with the template",
            "Restart Claude Code, invoke, observe",
        ], RED_BRIGHT),
        ("TRACK B", "Run Andorra's benchmark", [
            "B1   Activation queries",
            "B2   Real benchmark scenario",
            "B3   Write your own scenario (advanced)",
        ], CYAN),
    ]
    cw = 5.5; gap = 0.5
    total = cw * 2 + gap
    start_x = (SLIDE_W_IN - total) / 2
    for i, (label, name, items, accent) in enumerate(cols):
        x = Inches(start_x + i * (cw + gap)); y = Inches(2.3); w = Inches(cw); h = Inches(4.0)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.4), y + Inches(0.4), Inches(2), Inches(0.3),
                 label, size=10, color=TEXT_GRAY, bold=True)
        add_text(s, x + Inches(0.4), y + Inches(0.7), w - Inches(0.8), Inches(0.5),
                 name, size=20, color=TEXT_DARK, bold=True)
        for j, it in enumerate(items):
            add_text(s, x + Inches(0.4), y + Inches(1.6 + j * 0.5), w - Inches(0.8), Inches(0.4),
                     "›  " + it, size=12, color=accent if i == 0 else TEXT_BODY)
    add_text(s, Inches(0.5), Inches(6.7), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "Both tracks produce a real artefact.",
             size=11, color=TEXT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 23 — Andorra's Eval Suite
# ════════════════════════════════════════════════════════════════
def slide_23_eval_suite():
    s = new_slide()
    add_title(s, "Andorra's Eval Suite")
    add_subtitle(s, "Reproducible ground truth, not invented exercises.")
    stats = [
        ("21", "trigger queries", "10 should-trigger · 11 should-not"),
        ("7", "scenario directories", "6 in main benchmark + 1 extra"),
        ("53/53", "with-skill checks pass", "vs 48/53 without (90.6%)"),
    ]
    cw = 3.6; gap = 0.4
    total = cw * 3 + gap * 2
    start_x = (SLIDE_W_IN - total) / 2
    for i, (n, l, sub) in enumerate(stats):
        x = Inches(start_x + i * (cw + gap)); y = Inches(2.5); w = Inches(cw); h = Inches(2.8)
        add_card(s, x, y, w, h)
        add_text(s, x, y + Inches(0.6), w, Inches(1.0),
                 n, size=42, color=RED, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.7), w, Inches(0.4),
                 l, size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(2.1), w, Inches(0.4),
                 sub, size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.3), Inches(SLIDE_W_IN - 1), Inches(0.5),
             "Pasted into Claude Code. Graded against the published rubric. The numbers come from a JSON file you can read.",
             size=11, color=TEXT_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 24 — Key Takeaways
# ════════════════════════════════════════════════════════════════
def slide_24_takeaways():
    s = new_slide()
    add_title(s, "Key Takeaways")
    add_subtitle(s, "Slow productivity — skills force thinking before coding.")
    items = [
        "A skill is a markdown file that teaches an agent your methodology.",
        "Skills encode domain expertise, enforce guardrails, prevent common mistakes.",
        "Anatomy: frontmatter + steps + rules + gotchas.",
        "Start simple: 3–8 steps, 1–3 rules, 1–2 gotchas — refine over time.",
        "Share your skills: reusable, composable, platform-portable.",
    ]
    y = Inches(2.2)
    for i, t in enumerate(items):
        x = Inches(2); w = Inches(SLIDE_W_IN - 4); h = Inches(0.7)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.3), y + Inches(0.18), Inches(0.6), Inches(0.5),
                 str(i + 1), size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.0), y + Inches(0.2), w - Inches(1.3), Inches(0.4),
                 t, size=14, color=TEXT_DARK)
        y += Inches(0.85)
    add_text(s, Inches(0.5), Inches(6.7), Inches(SLIDE_W_IN - 1), Inches(0.4),
             "The value is in the guardrails: the things you would otherwise forget under time pressure.",
             size=12, color=EMERALD, bold=True, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 25 — Resources & Next Steps
# ════════════════════════════════════════════════════════════════
def slide_25_resources():
    s = new_slide()
    add_title(s, "Resources & Next Steps")
    add_subtitle(s, "Where to go from here.")
    resources = [
        ("Agent Skills Spec", "agentskills.io"),
        ("Bayesian skills", "github.com/Learning-Bayesian-Statistics/baygent-skills"),
        ("Software skills", "github.com/obra/superpowers"),
        ("Building Effective Agents", "anthropic.com/research/building-effective-agents"),
        ("Skill Creator", "/skill-creator (Claude Code plugin)"),
        ("Session 3 — Stubbs", "astrostubbs.github.io/GenAI-for-Scholarship"),
    ]
    cw = 5.5; gap = 0.3
    for i, (label, url) in enumerate(resources):
        col = i % 2; row = i // 2
        x = Inches((SLIDE_W_IN - cw * 2 - gap) / 2 + col * (cw + gap))
        y = Inches(2.3 + row * 1.0)
        w = Inches(cw); h = Inches(0.85)
        add_card(s, x, y, w, h)
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 "›", size=20, color=RED, bold=True)
        add_text(s, x + Inches(0.85), y + Inches(0.15), w - Inches(1.1), Inches(0.35),
                 label, size=13, color=TEXT_DARK, bold=True)
        add_text(s, x + Inches(0.85), y + Inches(0.5), w - Inches(1.1), Inches(0.3),
                 url, size=10, color=TEXT_GRAY, font=FONT_MONO)
    # CTA box
    bx = Inches((SLIDE_W_IN - 7) / 2); by = Inches(6.2); bw = Inches(7); bh = Inches(1.0)
    add_card(s, bx, by, bw, bh, fill=RED_TINT, border=RED_BRIGHT, border_w=1.0)
    add_text(s, bx, by + Inches(0.18), bw, Inches(0.4),
             "Refine. Install. Test. Share.",
             size=18, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, bx, by + Inches(0.6), bw, Inches(0.3),
             "The skill you started today is a v0.1. Iterate it on real work this week.",
             size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# Build all slides
# ════════════════════════════════════════════════════════════════
SLIDE_BUILDERS = [
    slide_01_title,
    slide_02_where_we_are,
    slide_03_landscape,
    slide_04_roadmap,
    slide_05_chatbot_vs_agent,
    slide_06_agent_loop,
    slide_07_four_components,
    slide_08_workflows_vs_agents,
    slide_09_claude_code_is,
    slide_10_problem,
    slide_11_what_is_skill,
    slide_12_why_skills_matter,
    slide_13_skills_vs_else,
    slide_14_skill_md_anatomy,
    slide_15_bayesian_case,
    slide_16_bayesian_results,
    slide_17_stress_test,
    slide_18_causal_inference,
    slide_19_more_skills,
    slide_20_practical_example,
    slide_21_your_research,
    slide_22_hands_on,
    slide_23_eval_suite,
    slide_24_takeaways,
    slide_25_resources,
]

for builder in SLIDE_BUILDERS:
    builder()

OUT = "/Users/cy/dev/hdsi/workshop_session4/slides_app.pptx"
prs.save(OUT)
print(f"Wrote {OUT}")
print(f"  {len(SLIDE_BUILDERS)} slides, 16:9 ({SLIDE_W_IN}\" × {SLIDE_H_IN}\")")
